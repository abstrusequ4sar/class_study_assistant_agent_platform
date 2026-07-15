"""生成简洁、可编辑的课程学习助手 Agent 平台汇报 PPT。"""
from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
ASSETS = ROOT / "assets" / "screenshots"
OUTPUT = ROOT / "课程学习助手Agent平台新版汇报.pptx"

NAVY = "1F4E79"
BLUE = "2F75B5"
LIGHT_BLUE = "EAF2F8"
PALE_BLUE = "F4F8FC"
GREEN = "70AD47"
LIGHT_GREEN = "EAF4E4"
ORANGE = "ED7D31"
LIGHT_ORANGE = "FCEDE3"
RED = "C00000"
LIGHT_RED = "FBEAEA"
DARK = "263238"
MID = "5F6B73"
LINE = "D9E2F3"
LIGHT = "F6F8FA"
WHITE = "FFFFFF"
FONT = "Microsoft YaHei"


def color(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def set_font(run, size=18, bold=False, color_hex=DARK):
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color(color_hex)


def add_text(
    slide,
    text,
    x,
    y,
    w,
    h,
    *,
    size=18,
    bold=False,
    color_hex=DARK,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin=0.04,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = valign
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.space_after = Pt(0)
    paragraph.line_spacing = 1.05
    run = paragraph.add_run()
    run.text = text
    set_font(run, size, bold, color_hex)
    return box


def add_bullets(slide, items, x, y, w, h, *, size=16, color_hex=DARK, gap=5):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.04)
    frame.margin_right = Inches(0.04)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = f"• {item}"
        paragraph.space_after = Pt(gap)
        paragraph.line_spacing = 1.08
        for run in paragraph.runs:
            set_font(run, size, False, color_hex)
    return box


def add_box(slide, x, y, w, h, *, fill=WHITE, line=LINE, radius=True, width=1):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color(fill)
    shape.line.color.rgb = color(line)
    shape.line.width = Pt(width)
    return shape


def add_title(slide, title, number, section=""):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.68))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color(NAVY)
    bar.line.fill.background()
    add_text(slide, title, 0.48, 0.10, 9.8, 0.44, size=24, bold=True, color_hex=WHITE)
    if section:
        add_text(slide, section, 10.15, 0.15, 2.55, 0.3, size=11, color_hex="DCE6F1", align=PP_ALIGN.RIGHT)
    add_text(slide, f"{number:02d}", 12.70, 7.08, 0.35, 0.22, size=9, color_hex=MID, align=PP_ALIGN.RIGHT)


def new_slide(prs, title, section=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = color(WHITE)
    add_title(slide, title, len(prs.slides), section)
    return slide


def add_section_card(slide, title, items, x, y, w, h, *, fill=PALE_BLUE, title_color=BLUE, body_size=15):
    add_box(slide, x, y, w, h, fill=fill, line=LINE)
    add_text(slide, title, x + 0.16, y + 0.11, w - 0.32, 0.30, size=16, bold=True, color_hex=title_color)
    add_bullets(slide, items, x + 0.16, y + 0.50, w - 0.32, h - 0.58, size=body_size, gap=4)


def add_picture_fit(slide, filename, x, y, w, h, *, caption=None):
    path = ASSETS / filename
    if not path.is_file():
        raise FileNotFoundError(path)
    add_box(slide, x, y, w, h, fill=WHITE, line="B4C7E7", radius=False, width=1.2)
    with Image.open(path) as image:
        image_ratio = image.width / image.height
    box_ratio = w / h
    if image_ratio >= box_ratio:
        pic_w = w - 0.08
        pic_h = pic_w / image_ratio
        pic_x = x + 0.04
        pic_y = y + (h - pic_h) / 2
    else:
        pic_h = h - 0.08
        pic_w = pic_h * image_ratio
        pic_x = x + (w - pic_w) / 2
        pic_y = y + 0.04
    slide.shapes.add_picture(str(path), Inches(pic_x), Inches(pic_y), Inches(pic_w), Inches(pic_h))
    if caption:
        add_text(slide, caption, x, y + h + 0.04, w, 0.25, size=10, color_hex=MID, align=PP_ALIGN.CENTER)


def add_arrow(slide, x1, y1, x2, y2, *, line_color=BLUE, width=2):
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    connector.line.color.rgb = color(line_color)
    connector.line.width = Pt(width)
    try:
        connector.line.end_arrowhead = True
    except (AttributeError, ValueError):
        pass
    return connector


def add_label_box(slide, text, x, y, w, h, *, fill=LIGHT_BLUE, line=BLUE, size=15, bold=False):
    add_box(slide, x, y, w, h, fill=fill, line=line)
    add_text(
        slide,
        text,
        x + 0.08,
        y + 0.05,
        w - 0.16,
        h - 0.10,
        size=size,
        bold=bold,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )


def feature_slide(prs, title, why, methods, benefits, image, caption, section="功能与实现"):
    slide = new_slide(prs, title, section)
    add_section_card(slide, "为什么需要", why, 0.48, 1.02, 4.25, 1.25, fill=LIGHT_ORANGE, title_color=ORANGE, body_size=14)
    add_section_card(slide, "实现方法", methods, 0.48, 2.43, 4.25, 2.65, fill=PALE_BLUE, title_color=BLUE, body_size=14)
    add_section_card(slide, "亮点与优势", benefits, 0.48, 5.24, 4.25, 1.42, fill=LIGHT_GREEN, title_color=GREEN, body_size=14)
    add_picture_fit(slide, image, 4.98, 1.10, 7.85, 5.42, caption=caption)
    return slide


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1 封面
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color(WHITE)
    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.28), Inches(7.5)).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = color(NAVY)
    slide.shapes[-1].line.fill.background()
    add_text(slide, "课程学习助手 Agent 平台", 0.82, 1.25, 11.6, 0.72, size=34, bold=True, color_hex=NAVY)
    add_text(slide, "面向多课程学习的资料检索、智能问答、学习规划与任务管理平台", 0.86, 2.15, 11.2, 0.55, size=20, color_hex=MID)
    labels = ["资料可检索", "Agent 可执行", "来源可跳转", "计划可落地"]
    for index, label in enumerate(labels):
        add_label_box(slide, label, 0.86 + index * 2.75, 3.10, 2.35, 0.62, fill=PALE_BLUE, line="9DC3E6", size=16, bold=True)
    add_text(slide, "大型程序设计实践汇报", 0.86, 5.22, 6.0, 0.40, size=20, bold=True, color_hex=BLUE)
    add_text(slide, "2026 年 7 月", 0.86, 5.78, 3.2, 0.32, size=15, color_hex=MID)
    add_text(slide, "汇报重点：为什么做、每项功能如何实现、实际优势是什么", 0.86, 6.63, 11.2, 0.34, size=13, color_hex=MID)

    # 2 为什么做
    slide = new_slide(prs, "为什么做：学习问题不是“缺一个聊天框”", "项目背景")
    pains = [
        ("01", "资料分散", "课件、笔记、作业和实验文档分布在不同位置，找资料成本高。"),
        ("02", "重点难找", "文档数量增加后，仅靠文件名不能快速定位知识片段。"),
        ("03", "计划不清", "“两周复习完”仍是宏观目标，缺少阶段与每日可执行安排。"),
        ("04", "任务易漏", "作业、实验、复习计划持续产生，只靠记忆难以跟踪。"),
    ]
    for index, (num, title, desc) in enumerate(pains):
        x = 0.58 + (index % 2) * 6.16
        y = 1.12 + (index // 2) * 2.40
        add_box(slide, x, y, 5.72, 2.02, fill=WHITE, line="B4C7E7")
        add_text(slide, num, x + 0.22, y + 0.20, 0.62, 0.50, size=25, bold=True, color_hex=BLUE)
        add_text(slide, title, x + 1.05, y + 0.20, 3.9, 0.42, size=21, bold=True, color_hex=NAVY)
        add_text(slide, desc, x + 1.05, y + 0.76, 4.20, 0.87, size=15, color_hex=MID)
    add_box(slide, 0.58, 6.05, 11.88, 0.68, fill=LIGHT_BLUE, line="9DC3E6")
    add_text(slide, "因此，需要的不只是回答问题，而是将资料、问答、计划和任务串成一个可持续执行的学习闭环。", 0.83, 6.20, 11.30, 0.30, size=16, bold=True, color_hex=NAVY, align=PP_ALIGN.CENTER)

    # 3 项目目标
    slide = new_slide(prs, "项目目标：建立以课程为中心的学习闭环", "项目背景")
    steps = ["资料入库", "内容检索", "Agent 推理", "计划拆解", "任务执行", "进度反馈"]
    fills = [LIGHT_BLUE, LIGHT_BLUE, LIGHT_GREEN, LIGHT_GREEN, LIGHT_ORANGE, PALE_BLUE]
    for index, step in enumerate(steps):
        x = 0.53 + index * 2.10
        add_label_box(slide, step, x, 2.10, 1.62, 0.82, fill=fills[index], line=BLUE if index < 2 else GREEN if index < 4 else ORANGE, size=15, bold=True)
        if index < len(steps) - 1:
            add_text(slide, "→", x + 1.68, 2.30, 0.34, 0.32, size=22, bold=True, color_hex=BLUE, align=PP_ALIGN.CENTER)
    add_section_card(slide, "对学生", ["从“找文件”变成“找答案”", "从“有目标”变成“今天做什么”"], 0.65, 3.58, 3.78, 1.78, fill=PALE_BLUE, title_color=BLUE, body_size=15)
    add_section_card(slide, "对系统", ["业务 CRUD 与 Agent 工具统一", "所有智能结果可存储、可追溯、可执行"], 4.78, 3.58, 3.78, 1.78, fill=LIGHT_GREEN, title_color=GREEN, body_size=15)
    add_section_card(slide, "对工程", ["外部 AI 不可用时基础功能仍可运行", "模型、数据库和检索服务可替换"], 8.91, 3.58, 3.78, 1.78, fill=LIGHT_ORANGE, title_color=ORANGE, body_size=15)
    add_text(slide, "最终目标：不让 AI 停留在“给建议”，而是把建议转化为可核查的答案和可落地的学习任务。", 1.05, 6.15, 11.25, 0.34, size=17, bold=True, color_hex=NAVY, align=PP_ALIGN.CENTER)

    # 4 需求全景
    slide = new_slide(prs, "需求完成全景", "需求分析")
    rows = [
        ("基本功能", "认证、课程、资料、Agent 对话、学习计划、待办、个人中心", "已完成", LIGHT_BLUE),
        ("后端能力", "8 张业务表、34 个 OpenAPI 操作、RESTful CRUD、文件存储", "已完成", PALE_BLUE),
        ("高级功能", "来源引用与跳转、知识点整理、智能任务拆解、多课程规划", "已完成", LIGHT_GREEN),
        ("工程增强", "SSE 流式输出、双协议模型、离线降级、日历提醒、XSS 消毒", "已完成", LIGHT_ORANGE),
    ]
    y = 1.18
    for label, detail, status, fill in rows:
        add_box(slide, 0.58, y, 12.15, 1.13, fill=fill, line="B4C7E7")
        add_text(slide, label, 0.82, y + 0.27, 1.55, 0.38, size=18, bold=True, color_hex=NAVY)
        add_text(slide, detail, 2.45, y + 0.20, 8.75, 0.62, size=15, color_hex=DARK, valign=MSO_ANCHOR.MIDDLE)
        add_label_box(slide, f"✓ {status}", 11.30, y + 0.27, 1.10, 0.48, fill=WHITE, line=GREEN, size=13, bold=True)
        y += 1.35
    add_text(slide, "需求对照的核心是：每个功能都有对应页面、API、数据结构和验证方式。", 0.74, 6.62, 11.8, 0.30, size=15, bold=True, color_hex=BLUE, align=PP_ALIGN.CENTER)

    # 5 总体架构
    slide = new_slide(prs, "系统总体架构", "架构设计")
    layers = [
        ("表现层", "Vue 3 / Element Plus / Pinia / Vue Router", LIGHT_BLUE, BLUE),
        ("API 层", "FastAPI 路由：auth · courses · materials · chat · plans · tasks", PALE_BLUE, NAVY),
        ("业务服务层", "Agent · LLM 统一客户端 · Retrieval · Extraction · Embedding · Security", LIGHT_GREEN, GREEN),
        ("数据与外部服务", "SQLite + 上传文件    |    Anthropic / OpenAI 兼容模型    |    Embedding 服务", LIGHT_ORANGE, ORANGE),
    ]
    y = 1.00
    for index, (name, detail, fill, line_color) in enumerate(layers):
        add_box(slide, 1.02, y, 11.25, 0.84, fill=fill, line=line_color, width=1.4)
        add_text(slide, name, 1.28, y + 0.09, 1.75, 0.50, size=17, bold=True, color_hex=line_color, valign=MSO_ANCHOR.MIDDLE)
        add_text(slide, detail, 3.25, y + 0.08, 8.55, 0.52, size=14, color_hex=DARK, valign=MSO_ANCHOR.MIDDLE)
        if index < len(layers) - 1:
            add_text(slide, "↓", 6.34, y + 0.84, 0.45, 0.23, size=16, bold=True, color_hex=BLUE, align=PP_ALIGN.CENTER)
        y += 1.15
    add_box(slide, 1.02, 5.72, 11.25, 0.92, fill=WHITE, line="B4C7E7")
    add_text(slide, "架构优势", 1.28, 5.89, 1.20, 0.32, size=16, bold=True, color_hex=NAVY)
    advantages = ["前后端分离", "模型/检索/数据库可替换", "API 与 Agent 工具复用权限逻辑"]
    for index, text in enumerate(advantages):
        add_label_box(slide, text, 2.68 + index * 3.02, 5.87, 2.74, 0.46, fill=PALE_BLUE if index < 2 else LIGHT_GREEN, line=BLUE if index < 2 else GREEN, size=12, bold=True)

    # 6 数据与闭环
    slide = new_slide(prs, "核心数据与学习闭环", "架构设计")
    entities = [
        ("User", 0.62, 1.25), ("Course", 2.55, 1.25), ("Material", 4.65, 1.25), ("Chunk", 6.75, 1.25),
        ("Conversation", 2.55, 2.60), ("Message + Citation", 4.65, 2.60), ("StudyPlan", 8.95, 1.25), ("Task", 10.65, 2.60),
    ]
    for name, x, y in entities:
        add_label_box(slide, name, x, y, 1.55 if name != "Message + Citation" else 2.05, 0.62, fill=PALE_BLUE, line=BLUE, size=13, bold=True)
    relationships = [
        (2.17, 1.56, 2.55, 1.56), (4.10, 1.56, 4.65, 1.56), (6.20, 1.56, 6.75, 1.56),
        (3.32, 1.87, 3.32, 2.60), (4.10, 2.91, 4.65, 2.91),
        (1.40, 1.87, 9.40, 1.87), (10.05, 1.80, 10.95, 2.60),
    ]
    for coords in relationships:
        add_arrow(slide, *coords, line_color="7F8C8D", width=1.4)
    add_text(slide, "8 张业务表把“回答”和“执行结果”都持久化，避免智能功能成为一次性输出。", 0.82, 3.55, 11.65, 0.42, size=15, color_hex=MID, align=PP_ALIGN.CENTER)
    flow = ["上传资料", "抽取/切片", "检索与引用", "Agent 决策", "计划或任务落库", "日历跟踪"]
    for index, text in enumerate(flow):
        x = 0.55 + index * 2.10
        add_label_box(slide, text, x, 4.48, 1.60, 0.75, fill=LIGHT_GREEN if index >= 2 else LIGHT_BLUE, line=GREEN if index >= 2 else BLUE, size=13, bold=True)
        if index < len(flow) - 1:
            add_text(slide, "→", x + 1.65, 4.68, 0.38, 0.28, size=18, bold=True, color_hex=BLUE, align=PP_ALIGN.CENTER)
    add_box(slide, 1.18, 5.78, 10.98, 0.80, fill=LIGHT_ORANGE, line=ORANGE)
    add_text(slide, "闭环价值：资料给 Agent 提供事实，Agent 生成的计划变成任务，任务状态又为后续安排提供反馈。", 1.45, 6.00, 10.45, 0.33, size=16, bold=True, color_hex=ORANGE, align=PP_ALIGN.CENTER)

    # 7 用户认证
    feature_slide(
        prs,
        "用户认证",
        ["学习资料、对话和计划属于个人数据，需要明确访问边界"],
        ["PBKDF2 加盐哈希保存密码", "JWT 签名与过期时间维护会话", "Vue Router 守卫拦截未登录页面", "Pinia 在页面刷新后通过 /auth/me 恢复用户信息"],
        ["所有业务接口复用当前用户依赖", "认证、路由与 API 错误处理保持统一"],
        "login.png",
        "系统统一登录入口",
    )

    # 8 个人中心
    feature_slide(
        prs,
        "个人中心",
        ["课程、计划、任务和对话分布在不同模块，需要统一概览入口"],
        ["并行获取课程、学习计划、未完成任务和对话数量", "展示最近对话并提供课程快捷跳转", "支持修改昵称和密码", "统计卡片可直接进入课程、计划和任务页"],
        ["把分散的学习数据汇总成个人仪表盘", "复用已有 RESTful API，不增加重复统计端点"],
        "profile.png",
        "个人中心统计、账号信息与最近对话",
    )

    # 8 课程管理
    feature_slide(
        prs,
        "课程管理",
        ["多门课程的资料、问答和任务需要清晰边界"],
        ["课程字段包含名称、简介、教师和学期", "FastAPI + SQLAlchemy 完成 CRUD", "前端以响应式卡片展示并提供问答快捷入口", "删除时处理资料、对话及计划任务关系"],
        ["以课程作为资料和 Agent 的隔离边界", "不同用户访问他人对象统一返回不存在"],
        "courses.png",
        "课程卡片同时展示核心信息与常用操作",
    )

    # 9 资料管理
    feature_slide(
        prs,
        "课程资料管理",
        ["课件、笔记、作业和实验文档需要统一归档和检索"],
        ["支持拖拽或点击选择，单批最多 10 个文件", "批量队列校验空文件、重复文件和 50MB 大小限制", "成功文件移出队列，失败文件保留重试", "原文件保存到上传目录，元数据写入 materials 表"],
        ["资料上传、分类与正文检索集中呈现", "批量有进度反馈，失败文件可直接重试"],
        "material_search.png",
        "拖拽批量上传、分类列表和正文检索",
    )

    # 10 解析与检索
    slide = new_slide(prs, "多格式解析与资料检索", "功能与实现")
    pipeline = [
        ("TXT / MD\nPDF / DOCX / PPTX", LIGHT_BLUE, BLUE),
        ("文本抽取\n页码/幻灯片标记", PALE_BLUE, BLUE),
        ("600 字切片\n100 字重叠", LIGHT_GREEN, GREEN),
        ("Embedding\n余弦相似度", LIGHT_GREEN, GREEN),
        ("语义检索\nTop-K 片段", LIGHT_ORANGE, ORANGE),
    ]
    for index, (text, fill, line_color) in enumerate(pipeline):
        x = 0.52 + index * 2.55
        add_label_box(slide, text, x, 1.48, 1.95, 1.05, fill=fill, line=line_color, size=14, bold=True)
        if index < len(pipeline) - 1:
            add_text(slide, "→", x + 2.00, 1.82, 0.40, 0.30, size=20, bold=True, color_hex=BLUE, align=PP_ALIGN.CENTER)
    add_section_card(slide, "为什么这样做", ["通用模型不知道学生自己的课程资料", "直接将整份文档放入上下文成本高且无法精确引用"], 0.62, 3.25, 3.72, 1.72, fill=LIGHT_ORANGE, title_color=ORANGE, body_size=14)
    add_section_card(slide, "如何容错", ["未配置 Embedding 或调用失败时自动改用关键词打分", "最多扫描单课程 3000 个切片，限制计算范围"], 4.75, 3.25, 3.72, 1.72, fill=PALE_BLUE, title_color=BLUE, body_size=14)
    add_section_card(slide, "亮点与优势", ["向量检索处理同义表达，关键词检索保证离线可用", "每个命中片段保留资料 ID 和切片 ID，为引用追溯打基础"], 8.88, 3.25, 3.82, 1.72, fill=LIGHT_GREEN, title_color=GREEN, body_size=14)
    add_box(slide, 1.10, 5.65, 11.05, 0.78, fill=LIGHT_BLUE, line="9DC3E6")
    add_text(slide, "结果：用户可以按“导数”、“进程调度”等知识词定位正文，Agent 也可用多组关键词自主检索。", 1.38, 5.88, 10.50, 0.30, size=16, bold=True, color_hex=NAVY, align=PP_ALIGN.CENTER)

    # 11 Agent
    feature_slide(
        prs,
        "完整工具驱动 Agent 对话",
        ["固定的“检索一次后作答”无法换词重试，也不能直接操作任务和计划"],
        ["向模型提供资料、课程、任务和计划共 10 个工具", "模型决定调用哪个工具、调多少轮、何时作答", "最多 8 轮循环，避免异常无限调用", "SSE 将 meta / tool / delta / done 事件实时展示"],
        ["流式与非流式接口复用同一 Agent 入口", "Anthropic 和 OpenAI 兼容协议共用业务工具集"],
        "agent.png",
        "Agent 在一轮对话中完成多次检索、查任务和创建任务",
    )

    # 13 来源展示
    feature_slide(
        prs,
        "资料来源展示",
        ["基于私有资料的回答需要让用户能够核查依据"],
        ["检索命中后以 chunk_id 去重并分配稳定编号", "模型看到的工具结果包含 [编号] 引用提示", "citations_json 保存资料 ID、名称和原文摘录", "回答正文显示编号，底部显示可悬停来源标签"],
        ["多轮检索中同一片段不会重复编号", "历史消息保留当时的引用快照"],
        "agent_sources_clickable.png",
        "Agent 回答正文、编号引用与参考资料标签",
        section="高级功能",
    )

    # 14 来源跳转
    feature_slide(
        prs,
        "资料来源点击跳转",
        ["只显示资料名仍需要用户自己在资料列表中再次查找"],
        ["来源标签使用 material_id 和引用编号构造路由查询参数", "点击或键盘 Enter 均可跳转到课程资料页", "资料页根据 material_id 定位对应行", "页面显示来源提示，并以蓝色背景和左边框高亮目标"],
        ["从回答到原始资料只需一次点击", "用户可继续下载原文或检索资料内容"],
        "source_jump.png",
        "跳转后显示引用编号、资料名并高亮目标资料",
        section="高级功能",
    )

    # 13 知识点
    feature_slide(
        prs,
        "知识点整理",
        ["资料已归档并不等于已形成可复习的知识结构"],
        ["读取全部切片并识别 Chap / Chapter / 第 N 章", "每批最多 50 个切片或约 30000 字符", "后端同时处理最多 3 批，按原位置组装", "SSE 实时返回批次进度和耗时"],
        ["兼顾全部章节覆盖与生成速度", "进度可见，最终顺序与全局 [编号] 保持稳定"],
        "knowledge_summary.png",
        "按资料章节与原文顺序生成完整复习提纲",
    )

    # 14 单课程计划
    feature_slide(
        prs,
        "单课程学习计划",
        ["“两周复习完”缺少阶段划分和每日工作量"],
        ["输入课程、学习目标、截止日期和每日时长", "返回 overview、stages 和 daily_tasks 结构化结果", "Pydantic 校验日期、时长和任务数量", "计划保存后自动批量创建待办"],
        ["生成的不是一段文字，而是可编辑、可跟踪的计划记录", "离线模式用确定性日期规则生成基础计划"],
        "plans.png",
        "学习目标被转换为阶段任务和每日待办",
    )

    # 15 任务拆解
    slide = new_slide(prs, "智能任务拆解", "高级功能")
    flow = [
        ("目标 + 期限 + 每日时长", LIGHT_BLUE, BLUE),
        ("结构化 JSON 生成", PALE_BLUE, BLUE),
        ("Schema 校验与日期归一化", LIGHT_GREEN, GREEN),
        ("保存 StudyPlan", LIGHT_GREEN, GREEN),
        ("批量生成 Task", LIGHT_ORANGE, ORANGE),
    ]
    for index, (text, fill, line_color) in enumerate(flow):
        x = 0.48 + index * 2.58
        add_label_box(slide, text, x, 1.55, 2.02, 0.92, fill=fill, line=line_color, size=14, bold=True)
        if index < len(flow) - 1:
            add_text(slide, "→", x + 2.06, 1.85, 0.38, 0.28, size=18, bold=True, color_hex=BLUE, align=PP_ALIGN.CENTER)
    add_section_card(slide, "为什么需要", ["学习目标往往粒度过大，不能直接执行", "单纯输出建议无法进入任务管理闭环"], 0.65, 3.10, 3.65, 1.70, fill=LIGHT_ORANGE, title_color=ORANGE, body_size=14)
    add_section_card(slide, "关键实现", ["阶段任务表达方向，每日任务表达可执行动作", "每条任务保留课程、计划和截止日期关联"], 4.78, 3.10, 3.65, 1.70, fill=PALE_BLUE, title_color=BLUE, body_size=14)
    add_section_card(slide, "边界与优势", ["单次最多生成 60 条待办，防止异常输出", "计划与任务分开存储，删除计划不会丢失已执行任务"], 8.91, 3.10, 3.77, 1.70, fill=LIGHT_GREEN, title_color=GREEN, body_size=14)
    add_box(slide, 1.38, 5.58, 10.55, 0.95, fill=LIGHT_BLUE, line=BLUE)
    add_text(slide, "亮点：Agent 在对话中也能直接调用 create_study_plan 或 create_task，将自然语言指令变成真实数据。", 1.68, 5.87, 9.95, 0.33, size=16, bold=True, color_hex=NAVY, align=PP_ALIGN.CENTER)

    # 16 多课程
    feature_slide(
        prs,
        "多课程综合学习规划",
        ["多门课程同时有截止日期时，单独规划容易造成时间冲突"],
        ["接收多组课程 ID、目标、期限和每日总时长", "将课程信息与各自目标同时提供给计划服务", "每条 daily_task 必须返回课程归属", "后端校验课程 ID 属于当前用户且在本次请求中"],
        ["保留课程归属，在任务页可分课程跟踪", "离线时按课程数和日期平均分配时间"],
        "plans.png",
        "高等数学与操作系统共享每日学习时间",
    )

    # 17 任务管理
    feature_slide(
        prs,
        "待办任务、日历与到期提醒",
        ["手动作业、计划生成任务和 Agent 创建任务需要统一跟踪"],
        ["任务支持创建、更新、完成、删除和课程筛选", "列表视图按紧急程度排序，日历视图按日期展示", "提醒接口把未完成任务分为已逾期、今日和三日内", "登录后全局提醒，任务页同步显示统计卡片"],
        ["三种任务来源共用一套数据和状态规则", "日历将“任务清单”变成“时间安排”"],
        "tasks_calendar.png",
        "到期统计、月度日历与任务完成状态",
    )

    # 18 API 与存储
    slide = new_slide(prs, "RESTful API 与数据存储", "后端服务")
    groups = [
        ("认证", "/auth/register\n/auth/login\n/auth/me"),
        ("课程", "/courses\n/courses/{id}"),
        ("资料", "/courses/{id}/materials\n/materials/{id}"),
        ("对话", "/conversations/{id}/messages\n/messages/stream"),
        ("计划", "/plans\n/plans/multi-course"),
        ("任务", "/tasks\n/tasks/reminders"),
    ]
    for index, (name, routes) in enumerate(groups):
        x = 0.48 + (index % 3) * 4.24
        y = 1.05 + (index // 3) * 1.48
        add_box(slide, x, y, 3.88, 1.18, fill=PALE_BLUE, line="9DC3E6")
        add_text(slide, name, x + 0.18, y + 0.12, 1.00, 0.32, size=16, bold=True, color_hex=NAVY)
        add_text(slide, routes, x + 1.15, y + 0.11, 2.50, 0.70, size=11, color_hex=MID, valign=MSO_ANCHOR.MIDDLE)
    add_section_card(slide, "实现方法", ["FastAPI 路由层处理身份依赖、参数校验和响应模型", "SQLAlchemy 2.0 表达对象关系，Pydantic v2 校验请求与响应", "上传文件与 SQLite 元数据分离存储", "普通消息与 SSE 消息共用完整 Agent 循环"], 0.64, 4.25, 5.76, 2.12, fill=PALE_BLUE, title_color=BLUE, body_size=14)
    add_section_card(slide, "亮点与优势", ["34 个 OpenAPI 操作自动生成 Swagger 文档", "DATABASE_URL 可将 SQLite 替换为其他数据库", "路由、Schema、ORM 和服务层分离，新增功能不必修改无关模块", "数据库与上传文件删除顺序考虑一致性"], 6.91, 4.25, 5.76, 2.12, fill=LIGHT_GREEN, title_color=GREEN, body_size=14)

    # 19 安全降级
    slide = new_slide(prs, "可用性、安全与降级设计", "工程质量")
    cards = [
        ("用户隔离", "每次课程、资料、对话、任务操作重新校验归属", BLUE, LIGHT_BLUE),
        ("密码与令牌", "PBKDF2 加盐哈希；JWT 签名、过期与登录失效处理", BLUE, PALE_BLUE),
        ("文件安全", "上传限制 50 MB，文件名净化，下载前再次检查所有权", GREEN, LIGHT_GREEN),
        ("前端渲染", "Markdown 经 DOMPurify 消毒，过滤恶意标签、事件和 URL", GREEN, LIGHT_GREEN),
        ("模型降级", "无 LLM 时问答返回检索片段，计划按确定性规则拆分", ORANGE, LIGHT_ORANGE),
        ("检索降级", "Embedding 未配置或失败时自动使用关键词检索", ORANGE, LIGHT_ORANGE),
    ]
    for index, (title, desc, line_color, fill) in enumerate(cards):
        x = 0.55 + (index % 3) * 4.20
        y = 1.10 + (index // 3) * 2.45
        add_box(slide, x, y, 3.85, 2.02, fill=fill, line=line_color)
        add_text(slide, title, x + 0.20, y + 0.18, 3.25, 0.36, size=18, bold=True, color_hex=line_color)
        add_text(slide, desc, x + 0.20, y + 0.68, 3.35, 0.92, size=14, color_hex=DARK)
    add_box(slide, 0.95, 6.12, 11.40, 0.64, fill=WHITE, line="B4C7E7")
    add_text(slide, "设计原则：AI 是可选增强能力，不是课程、资料和任务基础业务的单点依赖。", 1.18, 6.30, 10.90, 0.27, size=16, bold=True, color_hex=NAVY, align=PP_ALIGN.CENTER)

    # 20 测试
    slide = new_slide(prs, "测试覆盖与验证结果", "工程质量")
    metrics = [
        ("45", "后端 pytest", "含并发分批与 SSE 进度", BLUE),
        ("2", "前端 Vitest", "Markdown 安全渲染", GREEN),
        ("通过", "Vite 生产构建", "新增来源跳转后重新验证", ORANGE),
        ("0", "已知依赖漏洞", "npm audit 结果", NAVY),
    ]
    for index, (value, label, desc, accent) in enumerate(metrics):
        x = 0.55 + index * 3.13
        add_box(slide, x, 1.10, 2.87, 1.65, fill=WHITE, line=accent, width=1.5)
        add_text(slide, value, x + 0.12, 1.28, 2.60, 0.50, size=28, bold=True, color_hex=accent, align=PP_ALIGN.CENTER)
        add_text(slide, label, x + 0.12, 1.88, 2.60, 0.30, size=15, bold=True, color_hex=DARK, align=PP_ALIGN.CENTER)
        add_text(slide, desc, x + 0.18, 2.28, 2.48, 0.27, size=10, color_hex=MID, align=PP_ALIGN.CENTER)
    categories = [
        ("业务 CRUD", "认证、课程、资料、任务、计划"),
        ("资料处理", "DOCX/PPTX 抽取、切片、关键词与向量排序"),
        ("Agent 协议", "OpenAI tool_calls 与 Anthropic tool_use/tool_result"),
        ("数据安全", "跨用户访问、工具越权、XSS 消毒"),
        ("可用性", "离线问答、规划降级、SSE 中断与存档"),
    ]
    y = 3.30
    for index, (name, detail) in enumerate(categories):
        fill = PALE_BLUE if index % 2 == 0 else LIGHT_GREEN
        add_box(slide, 1.15, y, 11.05, 0.58, fill=fill, line=LINE)
        add_text(slide, name, 1.35, y + 0.12, 1.48, 0.27, size=14, bold=True, color_hex=NAVY)
        add_text(slide, detail, 3.05, y + 0.12, 8.65, 0.27, size=13, color_hex=DARK)
        y += 0.67

    # 21 亮点
    slide = new_slide(prs, "项目亮点与实际优势", "总结")
    highlights = [
        ("可执行 Agent", "不只回答，还可自主查资料、查任务、创建任务和保存计划。", BLUE),
        ("可核查来源", "稳定编号、片段摘录、历史快照和点击跳转共同降低不可追溯性。", GREEN),
        ("模型可替换", "Anthropic 和 OpenAI 兼容协议统一，通过环境变量切换模型与 Base URL。", ORANGE),
        ("在线/离线统一", "无密钥或外部服务失败时仍能演示完整基础业务和确定性降级。", NAVY),
        ("学习业务闭环", "资料入库、知识问答、计划拆解、任务提醒和进度反馈串成完整流程。", BLUE),
    ]
    y = 1.05
    for index, (title, detail, accent) in enumerate(highlights):
        add_box(slide, 0.72, y, 11.90, 0.92, fill=WHITE, line="B4C7E7")
        add_box(slide, 0.72, y, 0.12, 0.92, fill=accent, line=accent, radius=False)
        add_text(slide, f"{index + 1}", 1.02, y + 0.18, 0.38, 0.34, size=18, bold=True, color_hex=accent, align=PP_ALIGN.CENTER)
        add_text(slide, title, 1.58, y + 0.17, 2.25, 0.36, size=18, bold=True, color_hex=NAVY)
        add_text(slide, detail, 3.95, y + 0.14, 8.15, 0.52, size=14, color_hex=DARK, valign=MSO_ANCHOR.MIDDLE)
        y += 1.10

    # 22 总结与演示
    slide = new_slide(prs, "总结与现场演示路线", "总结")
    demo = ["登录并查看课程", "上传资料并检索", "Agent 多轮调用工具", "点击来源跳转", "生成多课程计划", "在日历中跟踪任务"]
    for index, text in enumerate(demo):
        x = 0.52 + index * 2.11
        add_label_box(slide, text, x, 1.35, 1.64, 0.90, fill=LIGHT_BLUE if index < 2 else LIGHT_GREEN if index < 4 else LIGHT_ORANGE, line=BLUE if index < 2 else GREEN if index < 4 else ORANGE, size=13, bold=True)
        if index < len(demo) - 1:
            add_text(slide, "→", x + 1.69, 1.66, 0.34, 0.27, size=18, bold=True, color_hex=BLUE, align=PP_ALIGN.CENTER)
    add_box(slide, 0.80, 3.05, 11.72, 1.15, fill=PALE_BLUE, line=BLUE)
    add_text(slide, "项目结论", 1.05, 3.25, 1.65, 0.35, size=20, bold=True, color_hex=NAVY)
    add_text(slide, "系统完成了从“管理学习资料”到“根据资料执行学习计划”的转变；Agent 的每个回答可核查，每个计划可落地，每个任务可跟踪。", 2.75, 3.22, 9.20, 0.58, size=16, color_hex=DARK, valign=MSO_ANCHOR.MIDDLE)
    add_section_card(slide, "三人协作", ["需求与基础后端", "Agent、检索、计划与测试", "前端、联调、文档与汇报"], 0.88, 4.72, 3.70, 1.55, fill=WHITE, title_color=BLUE, body_size=13)
    add_section_card(slide, "AI 伦理声明", ["人类负责基础规划与需求分析", "AI 仅负责细节代码填写", "所有结果由人类复核、测试并承担责任"], 4.82, 4.72, 3.70, 1.55, fill=LIGHT_GREEN, title_color=GREEN, body_size=13)
    add_section_card(slide, "谢谢，请指导", ["建议现场优先演示来源跳转和 Agent 创建任务", "问答时可回到架构、检索、安全与测试页"], 8.76, 4.72, 3.70, 1.55, fill=LIGHT_ORANGE, title_color=ORANGE, body_size=13)

    prs.core_properties.title = "课程学习助手 Agent 平台新版汇报"
    prs.core_properties.subject = "大型程序设计实践汇报"
    prs.core_properties.author = "课程学习助手 Agent 平台项目组"
    prs.core_properties.keywords = "Agent, RAG, 课程学习, FastAPI, Vue"
    prs.save(OUTPUT)
    print(OUTPUT)
    print(f"slides={len(prs.slides)}")


if __name__ == "__main__":
    build()
